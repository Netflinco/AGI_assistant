"""Generic 16:9 management-deck generator (Alibaba PuHuiTi default)."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from .policy import OfficePolicyError


DEFAULT_FONT = "Alibaba PuHuiTi"


def generate_pptx(spec: dict[str, Any], *, output_path: str | Path, brief: dict | None = None) -> Path:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
    except ImportError as exc:  # pragma: no cover
        raise OfficePolicyError("OFFICE_RUNTIME_DEPENDENCY_MISSING") from exc
    output = Path(output_path)
    output.parent.mkdir(mode=0o700, parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    for index, slide_spec in enumerate(spec["slides"], start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        _background(slide)
        _text(slide, 0.65, 0.45, 12.0, 0.7, slide_spec.get("title") or "", 28 if index > 1 else 34, True)
        body = slide_spec.get("body") or ""
        if body:
            _text(slide, 0.75, 1.45, 11.6, 3.8, body, 16, False)
        metrics = slide_spec.get("metrics") or []
        if metrics:
            _metrics(slide, metrics)
        if slide_spec.get("chart") == "bar" and metrics:
            chart_data = CategoryChartData()
            chart_data.categories = [str(item.get("label") or "")[:16] for item in metrics[:6]]
            values = [_numeric(item.get("value")) for item in metrics[:6]]
            chart_data.add_series("指标", values)
            slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(7.0), Inches(3.8), Inches(5.5), Inches(2.5), chart_data)
        source_ids = slide_spec.get("sources") or []
        if source_ids:
            citation_map = {item.get("evidence_id"): item for item in (brief or {}).get("citations", [])}
            labels = [f"{citation_map[item].get('publisher', '')} · {citation_map[item].get('canonical_url', '')}" for item in source_ids if item in citation_map]
            verification = "（含待核验信息）" if brief and brief.get("answer_status") == "PARTIALLY_VERIFIED" else ""
            _text(slide, 0.75, 6.55, 11.6, 0.55, f"截至：{brief.get('as_of') if brief else ''} {verification} 来源：" + "；".join(labels), 8, False)
        asset_sources = slide_spec.get("asset_sources") or []
        if asset_sources:
            _text(slide, 0.75, 6.15, 11.6, 0.28, "附件定位：" + "；".join(str(item) for item in asset_sources[:4]), 7, False)
        _text(slide, 12.25, 0.45, 0.45, 0.3, str(index), 8, False)
    presentation.save(output)
    return output


def _background(slide: Any) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = __import__("pptx.dml.color", fromlist=["RGBColor"]).RGBColor(248, 250, 252)


def _text(slide: Any, x: float, y: float, width: float, height: float, text: str, size: int, bold: bool) -> None:
    box = slide.shapes.add_textbox(__import__("pptx.util", fromlist=["Inches"]).Inches(x), __import__("pptx.util", fromlist=["Inches"]).Inches(y), __import__("pptx.util", fromlist=["Inches"]).Inches(width), __import__("pptx.util", fromlist=["Inches"]).Inches(height))
    frame = box.text_frame
    frame.clear()
    # The source document can be long even after bounded extraction.  Make
    # overflow prevention explicit in the generated OOXML instead of relying
    # on a viewer-specific default; G6 still verifies the rendered result.
    frame.word_wrap = True
    frame.auto_size = __import__("pptx.enum.text", fromlist=["MSO_AUTO_SIZE"]).MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    paragraph.text = str(text)
    for run in paragraph.runs:
        run.font.name = DEFAULT_FONT
        run.font.size = __import__("pptx.util", fromlist=["Pt"]).Pt(size)
        run.font.bold = bold
    paragraph.alignment = __import__("pptx.enum.text", fromlist=["PP_ALIGN"]).PP_ALIGN.LEFT


def _metrics(slide: Any, metrics: list[dict]) -> None:
    total = min(len(metrics), 4)
    for index, item in enumerate(metrics[:4]):
        x = 0.75 + index * 3.0
        _text(slide, x, 2.15, 2.6, 0.45, str(item.get("label") or "指标"), 12, False)
        _text(slide, x, 2.65, 2.6, 0.75, str(item.get("value") or "—"), 24, True)


def _numeric(value: object) -> float:
    try:
        return float(str(value).replace(",", "").replace("%", ""))
    except ValueError:
        return 0.0
