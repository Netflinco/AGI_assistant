"""Bounded, deterministic Office extraction.  No external requests are made."""

from __future__ import annotations

import csv
import io
import json
from pathlib import Path
from typing import Any

from .policy import OfficePolicyError, escape_csv_formula, validate_content_limits


def extract(filename: str, content: bytes) -> dict[str, Any]:
    suffix = Path(filename).suffix.lower()
    if suffix == ".xlsx":
        return _extract_xlsx(content)
    if suffix == ".docx":
        return _extract_docx(content)
    if suffix == ".csv":
        return _extract_csv(content)
    if suffix == ".pptx":
        return _extract_pptx(content)
    raise OfficePolicyError("OFFICE_UNSUPPORTED_TYPE")


def minimal_fragments(extraction: dict[str, Any], limit: int = 18) -> dict[str, Any]:
    """The only payload permitted to a model planner: no binary/full document."""
    kind = extraction.get("kind")
    if kind in {"xlsx", "csv"}:
        return {"kind": kind, "sheets": extraction.get("sheets", [])[:3], "metrics": extraction.get("metrics", [])[:limit],
                "source_refs": extraction.get("source_refs", [])[:limit], "data_classification": "INTERNAL", "purpose": "CREATE_MANAGEMENT_PPT"}
    return {"kind": kind, "headings": extraction.get("headings", [])[:limit], "paragraphs": extraction.get("paragraphs", [])[:limit],
            "source_refs": extraction.get("source_refs", [])[:limit], "data_classification": "INTERNAL", "purpose": "CREATE_MANAGEMENT_PPT"}


def _extract_xlsx(content: bytes) -> dict[str, Any]:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:  # pragma: no cover - surfaced in deployment readiness
        raise OfficePolicyError("OFFICE_RUNTIME_DEPENDENCY_MISSING") from exc
    book = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    formula_book = load_workbook(io.BytesIO(content), read_only=True, data_only=False)
    if len(book.sheetnames) > 20:
        raise OfficePolicyError("OFFICE_CONTENT_LIMIT_EXCEEDED")
    sheets, metrics, cells, formulas, source_refs = [], [], 0, [], []
    for sheet, formula_sheet in zip(book.worksheets, formula_book.worksheets):
        max_rows = min(sheet.max_row or 0, 100_001)
        if max_rows > 100_000:
            raise OfficePolicyError("OFFICE_CONTENT_LIMIT_EXCEEDED")
        preview: list[list[str]] = []
        for row_index, (row, formula_row) in enumerate(zip(sheet.iter_rows(values_only=True), formula_sheet.iter_rows(values_only=True)), start=1):
            row_values = ["" if value is None else str(value)[:120] for value in row]
            cells += sum(1 for value in row_values if value)
            if cells > 1_000_000:
                raise OfficePolicyError("OFFICE_CONTENT_LIMIT_EXCEEDED")
            if row_index <= 12:
                preview.append(row_values[:12])
            for column_index, formula_value in enumerate(formula_row, start=1):
                if isinstance(formula_value, str) and formula_value.startswith("="):
                    formulas.append({"source": f"{sheet.title}!{_column_name(column_index)}{row_index}", "formula": formula_value[:500]})
            if row_index > 100_000:
                break
        sheets.append({"name": sheet.title[:80], "rows": max_rows, "preview": preview})
        for preview_row_index, row in enumerate(preview[1:8], start=2):
            if len(row) >= 2 and row[0] and row[1]:
                source = f"{sheet.title}!A{preview_row_index}"
                metrics.append({"label": row[0], "value": row[1], "source": source})
                source_refs.append(source)
    validate_content_limits(sheets=len(sheets), rows=sum(item["rows"] for item in sheets), nonempty_cells=cells)
    return {"kind": "xlsx", "sheets": sheets, "metrics": metrics[:24], "formulas": formulas[:200], "source_refs": source_refs[:100], "nonempty_cells": cells}


def _extract_csv(content: bytes) -> dict[str, Any]:
    decoded = content.decode("utf-8-sig", errors="replace")
    reader = csv.reader(io.StringIO(decoded))
    preview, metrics, rows, nonempty_cells = [], [], 0, 0
    for row in reader:
        rows += 1
        if rows > 100_000:
            raise OfficePolicyError("OFFICE_CONTENT_LIMIT_EXCEEDED")
        nonempty_cells += sum(1 for value in row if value)
        if nonempty_cells > 1_000_000:
            raise OfficePolicyError("OFFICE_CONTENT_LIMIT_EXCEEDED")
        safe = [escape_csv_formula(value) for value in row[:20]]
        if rows <= 12:
            preview.append(safe)
    validate_content_limits(rows=rows, nonempty_cells=nonempty_cells)
    source_refs = []
    for preview_row_index, row in enumerate(preview[1:8], start=2):
        if len(row) >= 2:
            source = f"CSV!A{preview_row_index}"
            metrics.append({"label": row[0], "value": row[1], "source": source})
            source_refs.append(source)
    return {"kind": "csv", "sheets": [{"name": "CSV", "rows": rows, "preview": preview}], "metrics": metrics, "source_refs": source_refs, "nonempty_cells": nonempty_cells}


def _extract_docx(content: bytes) -> dict[str, Any]:
    try:
        from docx import Document
    except ImportError as exc:  # pragma: no cover
        raise OfficePolicyError("OFFICE_RUNTIME_DEPENDENCY_MISSING") from exc
    document = Document(io.BytesIO(content))
    paragraphs = [paragraph.text.strip()[:500] for paragraph in document.paragraphs if paragraph.text.strip()]
    page_breaks = sum(paragraph._p.xml.count("w:type=\"page\"") for paragraph in document.paragraphs)
    pages = page_breaks + 1
    if pages > 200:
        raise OfficePolicyError("OFFICE_CONTENT_LIMIT_EXCEEDED")
    headings = [text for text in paragraphs if len(text) < 120][:30]
    tables = []
    for table in document.tables[:12]:
        table_rows = [[cell.text.strip()[:120] for cell in row.cells[:12]] for row in table.rows[:12]]
        tables.append(table_rows)
    source_refs = [f"paragraph:{index}" for index, _text in enumerate(paragraphs[:120], start=1)]
    return {"kind": "docx", "headings": headings, "paragraphs": paragraphs[:120], "tables": tables, "source_refs": source_refs, "pages": pages}


def _extract_pptx(content: bytes) -> dict[str, Any]:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover
        raise OfficePolicyError("OFFICE_RUNTIME_DEPENDENCY_MISSING") from exc
    presentation = Presentation(io.BytesIO(content))
    if len(presentation.slides) > 100:
        raise OfficePolicyError("OFFICE_CONTENT_LIMIT_EXCEEDED")
    headings = []
    for slide in presentation.slides:
        text = " ".join(shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text).strip()
        if text:
            headings.append(text[:500])
    return {"kind": "pptx", "headings": headings, "paragraphs": headings, "slides": len(presentation.slides)}


def _column_name(number: int) -> str:
    """Excel column label without importing a second Office runtime helper."""
    letters = ""
    while number:
        number, remainder = divmod(number - 1, 26)
        letters = chr(65 + remainder) + letters
    return letters
