"""Deterministic G6 validation for generated Office artifacts."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from .policy import OfficePolicyError
from .generators import DEFAULT_FONT


def validate_pptx_structure(pptx_path: str | Path, spec: dict[str, Any]) -> None:
    """Reopen the deck and verify the approved specification survived writing."""
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover - deployment readiness failure
        raise OfficePolicyError("OFFICE_RUNTIME_DEPENDENCY_MISSING") from exc
    try:
        deck = Presentation(str(pptx_path))
    except Exception as exc:
        raise OfficePolicyError("OFFICE_STRUCTURE_INVALID") from exc
    expected = spec.get("slides") or []
    if len(deck.slides) != len(expected):
        raise OfficePolicyError("OFFICE_STRUCTURE_INVALID")
    for slide, slide_spec in zip(deck.slides, expected):
        all_text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text)
        title = str(slide_spec.get("title") or "")
        if title and title not in all_text:
            raise OfficePolicyError("OFFICE_STRUCTURE_INVALID")
        for metric in slide_spec.get("metrics") or []:
            value = str(metric.get("value") or "")
            if value and value not in all_text:
                raise OfficePolicyError("OFFICE_DATA_VALIDATION_FAILED")
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip() and run.font.name != DEFAULT_FONT:
                        raise OfficePolicyError("OFFICE_STYLE_INVALID")


def validate_rendered_preview(pdf_path: str | Path, png_path: str | Path, *, expected_pages: int) -> None:
    """Check that preview is a genuine non-empty rendering of every slide."""
    pdf = Path(pdf_path)
    png = Path(png_path)
    if not pdf.is_file() or pdf.stat().st_size < 16 or not png.is_file() or png.stat().st_size < 16:
        raise OfficePolicyError("OFFICE_RENDER_FAILED")
    try:
        from pypdf import PdfReader
        pages = len(PdfReader(str(pdf)).pages)
    except ImportError as exc:  # pragma: no cover
        raise OfficePolicyError("OFFICE_RUNTIME_DEPENDENCY_MISSING") from exc
    except Exception as exc:
        raise OfficePolicyError("OFFICE_RENDER_FAILED") from exc
    if pages != expected_pages:
        raise OfficePolicyError("OFFICE_RENDER_FAILED")
    try:
        from PIL import Image, ImageStat
        with Image.open(png) as image:
            image.verify()
        with Image.open(png) as image:
            rgb = image.convert("RGB")
            stat = ImageStat.Stat(rgb)
            # A completely monochrome page is a common symptom of failed
            # conversion.  This deliberately avoids OCR-based semantic claims.
            if max(stat.var) < 0.1:
                raise OfficePolicyError("OFFICE_RENDER_FAILED")
    except OfficePolicyError:
        raise
    except ImportError as exc:  # pragma: no cover
        raise OfficePolicyError("OFFICE_RUNTIME_DEPENDENCY_MISSING") from exc
    except Exception as exc:
        raise OfficePolicyError("OFFICE_RENDER_FAILED") from exc


def validate_rendered_pages(page_paths: list[str | Path], *, expected_pages: int) -> None:
    """Reject a render if any slide page is blank or missing.

    The downloadable PNG is a compact contact sheet, but visual validation
    must inspect every rasterised PDF page first.  Otherwise a good first page
    can hide a blank or missing page later in the deck.
    """
    if len(page_paths) != expected_pages:
        raise OfficePolicyError("OFFICE_RENDER_FAILED")
    try:
        from PIL import Image, ImageStat
        for page_path in page_paths:
            with Image.open(page_path) as image:
                image.verify()
            with Image.open(page_path) as image:
                rgb = image.convert("RGB")
                if max(ImageStat.Stat(rgb).var) < 0.1:
                    raise OfficePolicyError("OFFICE_RENDER_FAILED")
    except OfficePolicyError:
        raise
    except ImportError as exc:  # pragma: no cover
        raise OfficePolicyError("OFFICE_RUNTIME_DEPENDENCY_MISSING") from exc
    except Exception as exc:
        raise OfficePolicyError("OFFICE_RENDER_FAILED") from exc
