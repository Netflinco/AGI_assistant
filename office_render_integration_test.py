#!/usr/bin/env python3
"""Controlled runtime check for XLSX → PPTX → PDF/PNG rendering (GATE-606–609)."""

from __future__ import annotations

import io
import os
from pathlib import Path
import tempfile

with tempfile.TemporaryDirectory(prefix="agi-office-render-") as _tmp:
    os.environ["AGI_INSPECTION_DB"] = str(Path(_tmp) / "test.db")
    import server
    from office_agent.assets import OfficeAssetService
    from office_agent.jobs import OfficeJobService

    server.init_db(reset=True)
    conn = server.connect()
    from openpyxl import Workbook
    from docx import Document
    workbook = Workbook()
    workbook.active.title = "经营 KPI"
    workbook.active.append(["指标", "本周"])
    workbook.active.append(["成交额", 128])
    source = io.BytesIO(); workbook.save(source)
    assets = OfficeAssetService(conn, Path(_tmp) / "assets")
    asset = assets.create(tenant_id="tenant_jihu", user_id="u_admin", filename="kpi.xlsx", content=source.getvalue())
    service = OfficeJobService(conn, assets, Path(_tmp) / "artifacts")
    job = service.create_ppt_job(tenant_id="tenant_jihu", user_id="u_admin", conversation_id="conv", asset_ids=[asset["asset_id"]], title="KPI 管理层汇报")["job"]
    finished = service.run_job(job_id=job["job_id"], tenant_id="tenant_jihu", user_id="u_admin")
    assert finished["status"] == "SUCCEEDED", finished
    artifact = finished["artifacts"][0]
    pptx = service.get_artifact(version_id=artifact["version_id"], tenant_id="tenant_jihu", user_id="u_admin", kind="download")
    pdf = service.get_artifact(version_id=artifact["version_id"], tenant_id="tenant_jihu", user_id="u_admin", kind="preview")
    png = service.get_artifact(version_id=artifact["version_id"], tenant_id="tenant_jihu", user_id="u_admin", kind="preview_png")
    assert pptx and pdf and png and pptx[1].read_bytes()[:2] == b"PK" and pdf[1].read_bytes().startswith(b"%PDF") and png[1].read_bytes().startswith(b"\x89PNG")
    from pptx import Presentation
    deck = Presentation(str(pptx[1]))
    assert len(deck.slides) >= 2
    assert all(
        run.font.name == "Alibaba PuHuiTi"
        for slide in deck.slides for shape in slide.shapes if getattr(shape, "has_text_frame", False)
        for paragraph in shape.text_frame.paragraphs for run in paragraph.runs if run.text.strip()
    )
    # D1 covers both frozen input formats.  A Word report must follow the
    # same private extraction → Spec → PPTX → PDF/PNG delivery chain.
    document = Document()
    document.add_heading("经营周报", level=1)
    document.add_paragraph("本周成交额保持增长，重点跟进重点门店。")
    docx_raw = io.BytesIO(); document.save(docx_raw)
    docx_asset = assets.create(tenant_id="tenant_jihu", user_id="u_admin", filename="weekly.docx", content=docx_raw.getvalue())
    docx_job = service.create_ppt_job(tenant_id="tenant_jihu", user_id="u_admin", conversation_id="docx", asset_ids=[docx_asset["asset_id"]], title="Word 管理层汇报")["job"]
    docx_finished = service.run_job(job_id=docx_job["job_id"], tenant_id="tenant_jihu", user_id="u_admin")
    assert docx_finished["status"] == "SUCCEEDED", docx_finished
    docx_artifact = docx_finished["artifacts"][0]
    docx_pptx = service.get_artifact(version_id=docx_artifact["version_id"], tenant_id="tenant_jihu", user_id="u_admin", kind="download")
    docx_pdf = service.get_artifact(version_id=docx_artifact["version_id"], tenant_id="tenant_jihu", user_id="u_admin", kind="preview")
    docx_png = service.get_artifact(version_id=docx_artifact["version_id"], tenant_id="tenant_jihu", user_id="u_admin", kind="preview_png")
    assert docx_pptx and docx_pdf and docx_png
    docx_slide_count = len(Presentation(str(docx_pptx[1])).slides)
    assert docx_slide_count >= 2
    from PIL import Image
    with Image.open(docx_png[1]) as contact_sheet:
        # Two generated slides are packed into a two-column contact sheet,
        # proving every rendered page was kept for G6 inspection instead of
        # retaining only the first slide PNG.
        assert contact_sheet.width >= min(3, docx_slide_count) * 480
    conn.close()

print("PASS office render integration: PPTX opens and LibreOffice produced PDF + PNG preview")
