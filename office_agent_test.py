#!/usr/bin/env python3
"""P0 Office tests (asset boundary, immutable extraction, Spec, job and lifecycle)."""

from __future__ import annotations

import io
import json
import os
from pathlib import Path
import tempfile
from datetime import datetime, timezone
import zipfile
import warnings

from tests.fake_services import FakeModelGateway, FakeScanner, fake_preview, fake_renderer


class TestOfficeJobService:
    """Mixin factory to avoid a system LibreOffice dependency in unit tests."""
    @staticmethod
    def patch_preview(service):
        service._create_preview = fake_preview
        return service


with tempfile.TemporaryDirectory(prefix="agi-office-") as _tmp:
    os.environ["AGI_INSPECTION_DB"] = str(Path(_tmp) / "test.db")
    import server
    from agent_governance.policy_registry import set_feature
    from office_agent.assets import OfficeAssetService
    from office_agent.jobs import OfficeJobService
    from office_agent.policy import OfficePolicyError, inspect_asset, inspect_batch, research_brief_decision
    from office_agent.extraction import extract
    from office_agent.specs import deterministic_spec, validate_spec

    server.init_db(reset=True)
    conn = server.connect()
    now = lambda: datetime(2026, 8, 18, 12, tzinfo=timezone.utc)
    storage = Path(_tmp) / "assets"
    artifacts = Path(_tmp) / "artifacts"
    asset_service = OfficeAssetService(conn, storage, scanner=FakeScanner(), now=now)
    from openpyxl import Workbook
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "KPI"
    sheet.append(["指标", "数值"])
    sheet.append(["收入", 120])
    sheet.append(["合计", "=SUM(B2:B2)"])
    raw = io.BytesIO(); workbook.save(raw)
    asset = asset_service.create(tenant_id="tenant_jihu", user_id="u_admin", filename="weekly.xlsx", content=raw.getvalue())
    assert asset["sha256"] and asset_service.get(asset_id=asset["asset_id"], tenant_id="tenant_jihu", user_id="u_admin")
    assert asset_service.get(asset_id=asset["asset_id"], tenant_id="tenant_jihu", user_id="u_region") is None
    try:
        asset_service.create(tenant_id="tenant_jihu", user_id="u_admin", filename="fake.xlsx", content=b"not-a-zip")
        raise AssertionError("bad magic should be blocked")
    except OfficePolicyError as exc:
        assert exc.code == "OFFICE_MAGIC_MISMATCH"
    try:
        asset_service.create(tenant_id="tenant_jihu", user_id="u_admin", filename="secret.csv", content=b"API_KEY=demo_secret_123456789")
        raise AssertionError("strong sensitive content should be blocked")
    except OfficePolicyError as exc:
        assert exc.code == "OFFICE_STRONG_SENSITIVE_DATA"
    assert inspect_asset("weekly.xlsx", raw.getvalue(), declared_size=40 * 1024 * 1024).byte_size == 40 * 1024 * 1024
    try:
        inspect_asset("weekly.xlsx", raw.getvalue(), declared_size=40 * 1024 * 1024 + 1)
        raise AssertionError("size limit should be blocked")
    except OfficePolicyError as exc:
        assert exc.code == "OFFICE_FILE_TOO_LARGE"
    assert len(inspect_batch([("a.xlsx", raw.getvalue())] * 3, declared_sizes=[40 * 1024 * 1024] * 3)) == 3
    for bad_files, sizes, code in [
        ([("a.xlsx", raw.getvalue())] * 4, None, "OFFICE_BATCH_FILE_LIMIT_EXCEEDED"),
        ([("a.xlsx", raw.getvalue())] * 3, [40 * 1024 * 1024, 40 * 1024 * 1024, 40 * 1024 * 1024 + 1], "OFFICE_BATCH_SIZE_LIMIT_EXCEEDED"),
    ]:
        try:
            inspect_batch(bad_files, declared_sizes=sizes)
            raise AssertionError("batch policy should block")
        except OfficePolicyError as exc:
            assert exc.code == code
    macro = io.BytesIO(raw.getvalue())
    with zipfile.ZipFile(macro, "a") as archive:
        archive.writestr("xl/vbaProject.bin", b"not-executed")
    try:
        inspect_asset("macro.xlsx", macro.getvalue())
        raise AssertionError("macro OOXML should be blocked")
    except OfficePolicyError as exc:
        assert exc.code == "OFFICE_MACRO_FORBIDDEN"
    external = io.BytesIO(raw.getvalue())
    with warnings.catch_warnings():
        warnings.simplefilter("ignore", UserWarning)
        with zipfile.ZipFile(external, "a") as archive:
            archive.writestr("xl/_rels/workbook.xml.rels", b'<Relationships><Relationship Target="https://outside.example" TargetMode="External" /></Relationships>')
    try:
        inspect_asset("external.xlsx", external.getvalue())
        raise AssertionError("external OOXML relationship should be blocked")
    except OfficePolicyError as exc:
        assert exc.code == "OFFICE_EXTERNAL_REFERENCE_FORBIDDEN"
    compressed = io.BytesIO()
    with zipfile.ZipFile(compressed, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", b"x")
        archive.writestr("xl/workbook.xml", b"x")
        archive.writestr("xl/worksheets/sheet1.xml", b"x" * 1024 * 1024)
    try:
        inspect_asset("compressed.xlsx", compressed.getvalue())
        raise AssertionError("ratio-limited archive should be blocked")
    except OfficePolicyError as exc:
        assert exc.code == "OFFICE_ARCHIVE_LIMIT_EXCEEDED"
    # Normal DOCX style catalogues are highly compressible OOXML boilerplate;
    # they must not be confused with a user-controlled ZIP bomb.
    from docx import Document
    ordinary_docx = Document(); ordinary_docx.add_heading("周报", level=1); ordinary_docx.add_paragraph("正常 Word 输入")
    ordinary_docx_raw = io.BytesIO(); ordinary_docx.save(ordinary_docx_raw)
    assert inspect_asset("ordinary.docx", ordinary_docx_raw.getvalue()).detected_mime.endswith("wordprocessingml.document")
    too_many_slides = io.BytesIO()
    with zipfile.ZipFile(too_many_slides, "w", compression=zipfile.ZIP_STORED) as archive:
        archive.writestr("[Content_Types].xml", b"x")
        archive.writestr("ppt/presentation.xml", b"x")
        for index in range(1, 102):
            archive.writestr(f"ppt/slides/slide{index}.xml", b"x")
    try:
        inspect_asset("too-many-slides.pptx", too_many_slides.getvalue())
        raise AssertionError("PPT slide limit must block before worker scheduling")
    except OfficePolicyError as exc:
        assert exc.code == "OFFICE_CONTENT_LIMIT_EXCEEDED"
    too_many_pages = io.BytesIO()
    with zipfile.ZipFile(too_many_pages, "w", compression=zipfile.ZIP_STORED) as archive:
        archive.writestr("[Content_Types].xml", b"x")
        archive.writestr("word/document.xml", b"<w:document>" + b'<w:br w:type="page"/>' * 200 + b"</w:document>")
    try:
        inspect_asset("too-many-pages.docx", too_many_pages.getvalue())
        raise AssertionError("Word page limit must block before worker scheduling")
    except OfficePolicyError as exc:
        assert exc.code == "OFFICE_CONTENT_LIMIT_EXCEEDED"
    try:
        OfficeAssetService(conn, storage, scanner=FakeScanner("UNAVAILABLE"), now=now, require_scanner=True).create(tenant_id="tenant_jihu", user_id="u_admin", filename="unavailable.xlsx", content=raw.getvalue())
        raise AssertionError("unavailable production scanner should block")
    except OfficePolicyError as exc:
        assert exc.code == "OFFICE_VIRUS_SCAN_UNAVAILABLE"
    csv_extraction = extract("data.csv", b"name,value\n=cmd,1\n+sum,2\n")
    assert csv_extraction["sheets"][0]["preview"][1][0].startswith("'")
    from PIL import Image
    from office_agent.validation import validate_rendered_pages
    blank_page = Path(_tmp) / "blank.png"
    Image.new("RGB", (32, 32), "white").save(blank_page)
    try:
        validate_rendered_pages([blank_page], expected_pages=1)
        raise AssertionError("blank rendered pages must not pass G6")
    except OfficePolicyError as exc:
        assert exc.code == "OFFICE_RENDER_FAILED"
    xlsx_extraction = extract("weekly.xlsx", raw.getvalue())
    assert any(item["formula"] == "=SUM(B2:B2)" for item in xlsx_extraction["formulas"])
    assert research_brief_decision({"brief_id": "b"}, enabled=True, same_owner=True).reason_code == "RESEARCH_BRIEF_INVALID"
    model = FakeModelGateway()
    set_feature(conn, "tenant_jihu", "office_model_processing_enabled", True, server.now_iso())
    service = TestOfficeJobService.patch_preview(OfficeJobService(conn, asset_service, artifacts, model_gateway=model, renderer=fake_renderer, now=now))
    template_block = service.create_ppt_job(tenant_id="tenant_jihu", user_id="u_admin", conversation_id="template", asset_ids=[asset["asset_id"]], title="未授权模板", template_id="template_unapproved")
    assert template_block["status"] == "BLOCKED" and template_block["reason_code"] == "OFFICE_TEMPLATE_NOT_FOUND" and not model.calls
    created = service.create_ppt_job(tenant_id="tenant_jihu", user_id="u_admin", conversation_id="conv", asset_ids=[asset["asset_id"]], title="周报管理层 PPT")
    job_id = created["job"]["job_id"]
    assert service.create_ppt_job(tenant_id="tenant_jihu", user_id="u_admin", conversation_id="conv", asset_ids=[asset["asset_id"]], title="周报管理层 PPT")["deduped"]
    finished = service.run_job(job_id=job_id, tenant_id="tenant_jihu", user_id="u_admin")
    assert finished["status"] == "SUCCEEDED" and model.calls and finished["artifacts"], finished
    assert "storage_key" not in finished["artifacts"][0] and "source_refs_json" not in finished["artifacts"][0]
    found = service.get_artifact(version_id=finished["artifacts"][0]["version_id"], tenant_id="tenant_jihu", user_id="u_admin", kind="download")
    assert found and found[1].is_file()
    assert service.get_artifact(version_id=finished["artifacts"][0]["version_id"], tenant_id="tenant_jihu", user_id="u_region", kind="download") is None
    # Closing the approved-model policy must switch to deterministic planning
    # without sending a second document fragment to the model gateway.
    model_calls_before_disabled = len(model.calls)
    set_feature(conn, "tenant_jihu", "office_model_processing_enabled", False, server.now_iso())
    deterministic_job = service.create_ppt_job(tenant_id="tenant_jihu", user_id="u_admin", conversation_id="deterministic", asset_ids=[asset["asset_id"]], title="确定性 PPT")
    assert service.run_job(job_id=deterministic_job["job"]["job_id"], tenant_id="tenant_jihu", user_id="u_admin")["status"] == "SUCCEEDED"
    assert len(model.calls) == model_calls_before_disabled
    set_feature(conn, "tenant_jihu", "office_model_processing_enabled", True, server.now_iso())
    try:
        validate_spec({"schema_version": "SlideSpec.v1", "slides": [{"layout": "summary", "title": "x", "body": "x", "path": "/tmp/x", "sources": []}]})
        raise AssertionError("unsafe Spec should fail")
    except OfficePolicyError as exc:
        assert exc.code == "OFFICE_SPEC_INVALID"
    partial_brief = {
        "brief_id": "brief_partial", "producer_run_id": "res_partial",
        "owner": {"tenant_id": "tenant_jihu", "user_id": "u_admin"},
        "answer_status": "PARTIALLY_VERIFIED", "as_of": "2026-08-18T12:00:00+00:00", "freshness": "day", "topic": "政策",
        "claims": [{"claim_id": "claim_partial", "text": "政策仍待核验", "claim_status": "PARTIALLY_VERIFIED", "confidence": 0.5, "evidence_ids": ["ev_partial"]}],
        "citations": [{"evidence_id": "ev_partial", "title": "官方公告", "canonical_url": "https://official.example.com/policy", "publisher": "官方", "published_at": "2026-08-18T10:00:00+00:00", "fetched_at": "2026-08-18T12:00:00+00:00", "source_tier": "OFFICIAL"}],
        "limitations": ["信息截至 2026-08-18，待核验"], "content_hash": "partial_hash", "policy_version": "p0-2026-08-18", "expires_at": "2026-10-17T12:00:00+00:00",
    }
    conn.execute("INSERT INTO research_briefs VALUES (?,?,?,?,?,?,?,?,?)", ("brief_partial", "res_partial", "tenant_jihu", "u_admin", "PARTIALLY_VERIFIED", "partial_hash", json.dumps(partial_brief, ensure_ascii=False), partial_brief["expires_at"], now().isoformat()))
    set_feature(conn, "tenant_jihu", "research_to_office_enabled", False, server.now_iso())
    collaboration_disabled = service.create_ppt_job(tenant_id="tenant_jihu", user_id="u_admin", conversation_id="partial-off", brief_id="brief_partial", title="政策管理层 PPT")
    assert collaboration_disabled["status"] == "BLOCKED" and collaboration_disabled["reason_code"] == "RESEARCH_TO_OFFICE_DISABLED"
    set_feature(conn, "tenant_jihu", "research_to_office_enabled", True, server.now_iso())
    bad_status = {**partial_brief, "brief_id": "brief_bad_status", "answer_status": "NO_AUTHORITATIVE_SOURCE"}
    assert research_brief_decision(bad_status, enabled=True, same_owner=True).reason_code == "RESEARCH_BRIEF_STATUS_BLOCKED"
    partial_spec = deterministic_spec("政策管理层 PPT", {}, partial_brief)
    research_slide = next(slide for slide in partial_spec["slides"] if slide["title"] == "公开信息核验")
    sources_slide = next(slide for slide in partial_spec["slides"] if slide["layout"] == "sources")
    assert "待核验" in research_slide["body"] and "待核验" in sources_slide["title"]
    validate_spec(partial_spec, brief=partial_brief)
    uncited = deterministic_spec("政策管理层 PPT", {}, partial_brief)
    next(slide for slide in uncited["slides"] if slide["title"] == "公开信息核验")["sources"] = []
    next(slide for slide in uncited["slides"] if slide["layout"] == "sources")["sources"] = []
    try:
        validate_spec(uncited, brief=partial_brief)
        raise AssertionError("research claims without evidence must not generate a fact deck")
    except OfficePolicyError as exc:
        assert exc.code == "OFFICE_SPEC_UNCITED_RESEARCH"
    partial_job = service.create_ppt_job(tenant_id="tenant_jihu", user_id="u_admin", conversation_id="partial", brief_id="brief_partial", title="政策管理层 PPT")
    partial_finished = service.run_job(job_id=partial_job["job"]["job_id"], tenant_id="tenant_jihu", user_id="u_admin")
    assert partial_finished["status"] == "SUCCEEDED"
    partial_file = service.get_artifact(version_id=partial_finished["artifacts"][0]["version_id"], tenant_id="tenant_jihu", user_id="u_admin", kind="download")[1]
    from pptx import Presentation
    deck_text = "\n".join(shape.text for slide in Presentation(str(partial_file)).slides for shape in slide.shapes if hasattr(shape, "text"))
    assert "待核验" in deck_text and "截至" in deck_text
    blocked_action = service.create_ppt_job(tenant_id="tenant_jihu", user_id="u_admin", conversation_id="conv", asset_ids=[asset["asset_id"]], title="外部共享周报")
    assert blocked_action["reason_code"] == "OFFICE_ACTION_NOT_ALLOWED"
    retry_service = TestOfficeJobService.patch_preview(OfficeJobService(
        conn, asset_service, artifacts, renderer=lambda *_args, **_kwargs: (_ for _ in ()).throw(OfficePolicyError("OFFICE_RENDER_TIMEOUT")), now=now,
    ))
    retry_created = retry_service.create_ppt_job(tenant_id="tenant_jihu", user_id="u_admin", conversation_id="retry", asset_ids=[asset["asset_id"]], title="渲染重试样例")
    retry_job_id = retry_created["job"]["job_id"]
    failed = retry_service.run_job(job_id=retry_job_id, tenant_id="tenant_jihu", user_id="u_admin")
    assert failed["status"] == "FAILED" and failed["error_code"] == "OFFICE_RENDER_TIMEOUT"
    requeued = retry_service.retry_job(job_id=retry_job_id, tenant_id="tenant_jihu", user_id="u_admin")
    assert requeued and requeued["status"] == "QUEUED"
    retry_service.renderer = fake_renderer
    assert retry_service.run_job(job_id=retry_job_id, tenant_id="tenant_jihu", user_id="u_admin")["status"] == "SUCCEEDED"
    canceled = service.create_ppt_job(tenant_id="tenant_jihu", user_id="u_admin", conversation_id="cancel", asset_ids=[asset["asset_id"]], title="取消样例")
    assert service.cancel_job(job_id=canceled["job"]["job_id"], tenant_id="tenant_jihu", user_id="u_admin")["status"] == "CANCELED"
    # Cancellation that arrives during generation must win over a stale worker
    # update and must not leave a downloadable artifact behind.
    from office_agent.generators import generate_pptx
    cancellation_holder = {}
    def cancel_during_generation(spec, *, output_path, brief):
        generated = generate_pptx(spec, output_path=output_path, brief=brief)
        cancellation_holder["service"].cancel_job(job_id=cancellation_holder["job_id"], tenant_id="tenant_jihu", user_id="u_admin")
        return generated
    cancel_service = TestOfficeJobService.patch_preview(OfficeJobService(conn, asset_service, artifacts, generator=cancel_during_generation, renderer=fake_renderer, now=now))
    cancellation_holder["service"] = cancel_service
    in_flight = cancel_service.create_ppt_job(tenant_id="tenant_jihu", user_id="u_admin", conversation_id="cancel-running", asset_ids=[asset["asset_id"]], title="运行中取消样例")["job"]
    cancellation_holder["job_id"] = in_flight["job_id"]
    canceled_result = cancel_service.run_job(job_id=in_flight["job_id"], tenant_id="tenant_jihu", user_id="u_admin")
    assert canceled_result["status"] == "CANCELED" and canceled_result["artifacts"] == []
    assert not conn.execute("SELECT 1 FROM office_artifact_versions WHERE job_id=? AND status='SUCCEEDED'", (in_flight["job_id"],)).fetchone()
    clock = {"value": now()}
    expiring_assets = OfficeAssetService(conn, storage, scanner=FakeScanner(), now=lambda: clock["value"])
    cleanup_service = OfficeJobService(conn, expiring_assets, artifacts, now=lambda: clock["value"])
    clock["value"] = datetime(2026, 9, 18, 12, tzinfo=timezone.utc)
    assert cleanup_service.cleanup_expired() >= 3
    assert expiring_assets.get(asset_id=asset["asset_id"], tenant_id="tenant_jihu", user_id="u_admin") is None
    assert cleanup_service.get_artifact(version_id=finished["artifacts"][0]["version_id"], tenant_id="tenant_jihu", user_id="u_admin", kind="download") is None
    conn.close()

print("PASS office tests: asset gates, DLP, ACL, extraction, minimal-model input, Spec and private artifact")
